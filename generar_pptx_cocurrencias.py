"""
Genera 3 láminas PowerPoint con coocurrencias por períodos históricos.
Solo incluye documentos con tema CAMBIO CLIMÁTICO (excluye los 14 institucionales).

Ejecutar:
    python3 generar_pptx_cocurrencias.py

Salida: coocurrencias_periodos.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import networkx as nx
from adjustText import adjust_text
from collections import Counter
from itertools import combinations
import io

from datos import (
    cargar_datos,
    cargar_mapa_clusters,
    _mtime_archivo,
    ARCHIVO_DATOS,
    ARCHIVO_CLUSTERS,
    filtrar_cambio_climatico,
    filtrar_solo_sustantivas,
    filtrar_excluir_boletines,
    _excluir_documentos,
    color_tema,
    normalizar_nombre_tema,
    temas_sin_cambio_climatico,
)

CLUSTER_LEGEND = [
    ("#c0392b", "Desarrollo económico"),
    ("#2980b9", "Desarrollo social"),
    ("#27ae60", "Sustentabilidad ambiental"),
    ("#8e44ad", "Desarrollo productivo e innovación"),
    ("#d4ac0d", "Institucionalidad y temas transversales"),
]


def calcular_coocurrencias(df, top_temas=40, peso_min=2):
    freq = Counter()
    aristas = Counter()

    for valor in df["cepal.topicSpa"]:
        temas = [normalizar_nombre_tema(t) for t in temas_sin_cambio_climatico(valor)]
        for t in temas:
            freq[t] += 1
        if len(temas) < 2:
            continue
        for a, b in combinations(sorted(set(temas)), 2):
            aristas[(a, b)] += 1

    temas_top = {t for t, _ in freq.most_common(top_temas)}
    aristas_filtradas = Counter({
        (a, b): w
        for (a, b), w in aristas.items()
        if w >= peso_min and a in temas_top and b in temas_top
    })
    nodos = {a for a, _ in aristas_filtradas} | {b for _, b in aristas_filtradas}
    freq_filtrada = Counter({t: freq[t] for t in nodos})

    return aristas_filtradas, freq_filtrada


def strength_nodos(aristas):
    s = {}
    for (a, b), w in aristas.items():
        s[a] = s.get(a, 0.0) + w
        s[b] = s.get(b, 0.0) + w
    return s


def generar_grafo_coocurrencias(df, mapa, periodo_label, top_temas=30, peso_min=2):
    aristas, freq = calcular_coocurrencias(df, top_temas, peso_min)

    if not aristas:
        return None

    G = nx.Graph()
    for (a, b), w in aristas.items():
        G.add_edge(a, b, weight=w)

    strength = strength_nodos(aristas)

    fig, ax = plt.subplots(figsize=(13, 9.5))

    pos = nx.spring_layout(G, weight="weight", k=3.5, iterations=150, seed=42)

    node_sizes = [max(strength.get(n, 1) * 60, 120) for n in G.nodes()]

    node_colors = []
    for n in G.nodes():
        color_hex = color_tema(n, mapa)
        r = int(color_hex[1:3], 16) / 255
        g = int(color_hex[3:5], 16) / 255
        b = int(color_hex[5:7], 16) / 255
        node_colors.append((r, g, b))

    max_w = max(aristas.values()) if aristas else 1
    edge_widths = [0.5 + (w / max_w) * 3.5 for w in aristas.values()]

    nx.draw_networkx_edges(G, pos, width=edge_widths, alpha=0.35, edge_color="#888888", ax=ax)
    nx.draw_networkx_nodes(G, pos, node_size=node_sizes, node_color=node_colors,
                           alpha=0.85, edgecolors="white", linewidths=1.5, ax=ax)

    top_n = min(25, len(G.nodes()))
    top_nodes_list = sorted(strength.items(), key=lambda x: -x[1])[:top_n]
    top_set = {n for n, _ in top_nodes_list}

    x_coords = [pos[n][0] for n in G.nodes() if n in top_set]
    y_coords = [pos[n][1] for n in G.nodes() if n in top_set]

    sorted_strengths = sorted(strength.values(), reverse=True)
    threshold = sorted_strengths[min(4, len(sorted_strengths) - 1)] if sorted_strengths else 0

    texts = []
    for n in G.nodes():
        if n in top_set:
            fs = 8 if strength.get(n, 0) >= threshold else 6
            t = ax.text(pos[n][0], pos[n][1], n, fontsize=fs, fontweight="bold",
                        ha="center", va="center", color="#1a1a1a")
            texts.append(t)

    if texts:
        adjust_text(
            texts,
            x=x_coords,
            y=y_coords,
            force_text=(0.8, 0.8),
            force_points=(0.4, 0.4),
            expand=(1.3, 1.5),
            max_move=(15, 15),
            arrowprops=dict(arrowstyle="-", color="#aaaaaa", lw=0.5),
            ax=ax,
        )

    ax.set_title("Período: " + periodo_label, fontsize=13, fontweight="bold", pad=12)
    ax.axis("off")
    fig.subplots_adjust(left=0.02, right=0.98, top=0.95, bottom=0.02)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def crear_pptx(df_completo, mapa):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    periodos = [
        ("2015-2018", 2015, 2018),
        ("2019-2022", 2019, 2022),
        ("2023-2026", 2023, 2026),
    ]

    slide_layout = prs.slide_layouts[6]

    for periodo_label, anio_inicio, anio_fin in periodos:
        serie_anio = pd.to_numeric(df_completo["dc.year"], errors="coerce")
        mask = (serie_anio >= anio_inicio) & (serie_anio <= anio_fin)
        df_periodo = df_completo[mask]

        slide = prs.slides.add_slide(slide_layout)

        tx_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5))
        p = tx_title.text_frame.paragraphs[0]
        p.text = "¿Qué topics coocurren junto a cambio climático? – " + periodo_label
        p.font.size = Pt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        tx_sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.3))
        p2 = tx_sub.text_frame.paragraphs[0]
        p2.text = str(len(df_periodo)) + " publicaciones sobre cambio climático de la CEPAL"
        p2.font.size = Pt(11)
        p2.alignment = PP_ALIGN.CENTER
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        tx_nota = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.3), Inches(0.4))
        tf_nota = tx_nota.text_frame
        tf_nota.word_wrap = True
        p3 = tf_nota.paragraphs[0]
        p3.text = ("Cada nodo es un topic presente en las publicaciones. El tamaño refleja su frecuencia. "
                    "Los enlaces indican coocurrencia en un mismo documento; el grosor es proporcional "
                    "al número de veces que aparecen juntos.")
        p3.font.size = Pt(9)
        p3.font.italic = True
        p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p3.alignment = PP_ALIGN.CENTER

        img_buf = generar_grafo_coocurrencias(df_periodo, mapa, periodo_label)
        if img_buf:
            slide.shapes.add_picture(img_buf, Inches(0.8), Inches(1.45), Inches(9.5), Inches(5.0))
        else:
            tx_empty = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
            p_empty = tx_empty.text_frame.paragraphs[0]
            p_empty.text = "No hay datos suficientes para el período " + periodo_label
            p_empty.font.size = Pt(16)
            p_empty.alignment = PP_ALIGN.CENTER

        tx_leg_title = slide.shapes.add_textbox(Inches(10.6), Inches(1.5), Inches(2.5), Inches(0.3))
        p_lt = tx_leg_title.text_frame.paragraphs[0]
        p_lt.text = "Clusters"
        p_lt.font.size = Pt(10)
        p_lt.font.bold = True

        for i, (hex_color, nombre) in enumerate(CLUSTER_LEGEND):
            y = 1.9 + i * 0.35
            r = int(hex_color[1:3], 16)
            g = int(hex_color[3:5], 16)
            b = int(hex_color[5:7], 16)

            tx_color = slide.shapes.add_textbox(Inches(10.6), Inches(y), Inches(0.3), Inches(0.25))
            tf_c = tx_color.text_frame
            tf_c.paragraphs[0].text = "■"
            tf_c.paragraphs[0].font.size = Pt(14)
            tf_c.paragraphs[0].font.color.rgb = RGBColor(r, g, b)

            tx_name = slide.shapes.add_textbox(Inches(10.95), Inches(y), Inches(2.2), Inches(0.25))
            tf_n = tx_name.text_frame
            tf_n.paragraphs[0].text = nombre
            tf_n.paragraphs[0].font.size = Pt(8)
            tf_n.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        tx_src = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3))
        p4 = tx_src.text_frame.paragraphs[0]
        p4.text = "Fuente: CEPAL – Corpus de publicaciones sobre cambio climático"
        p4.font.size = Pt(8)
        p4.font.italic = True
        p4.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p4.alignment = PP_ALIGN.RIGHT

    archivo_salida = "coocurrencias_periodos.pptx"
    prs.save(archivo_salida)
    return archivo_salida


def main():
    print("Cargando datos...")
    df = cargar_datos(_mtime=_mtime_archivo(ARCHIVO_DATOS))
    df = filtrar_cambio_climatico(df)
    df = filtrar_solo_sustantivas(df)
    df = filtrar_excluir_boletines(df)
    df = _excluir_documentos(df)
    mapa = cargar_mapa_clusters(_mtime=_mtime_archivo(ARCHIVO_CLUSTERS))

    print("Documentos de cambio climático: " + str(len(df)))

    print("Generando PowerPoint...")
    archivo = crear_pptx(df, mapa)

    print("PowerPoint generado: " + archivo)
    print("Láminas creadas:")
    print("  1. Coocurrencias 2015-2018")
    print("  2. Coocurrencias 2019-2022")
    print("  3. Coocurrencias 2023-2026")


if __name__ == "__main__":
    main()
